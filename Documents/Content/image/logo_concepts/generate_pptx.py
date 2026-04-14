"""Generate Strata Forge Brand Identity PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

ICON_PATH = os.path.join(os.path.dirname(__file__), "strata_forge_icon.png")
OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "Strata_Forge_Brand_Identity.pptx")

# Brand colors
ABYSS = RGBColor(0x0A, 0x0A, 0x12)
TEAL = RGBColor(0x2A, 0x6B, 0x6B)
TEAL_LIGHT = RGBColor(0x3A, 0x8A, 0x8A)
ORANGE = RGBColor(0xE8, 0x78, 0x30)
ORANGE_DEEP = RGBColor(0xC4, 0x5A, 0x2C)
INCANDESCENT = RGBColor(0xE8, 0xE0, 0xD0)
WARM_WHITE = RGBColor(0xFF, 0xF5, 0xE0)
CONCRETE = RGBColor(0x77, 0x77, 0x77)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
LIGHT_BG = RGBColor(0xF0, 0xED, 0xE6)
TEAL_BG = RGBColor(0x14, 0x2E, 0x2E)
WARM_BG = RGBColor(0x1A, 0x12, 0x10)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size, color, bold=False, font_name='Cinzel', alignment=PP_ALIGN.CENTER, spacing=0):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if spacing > 0:
        p.font._element.attrib['{http://schemas.openxmlformats.org/drawingml/2006/main}spc'] = str(spacing)
    return txBox

def add_icon(slide, left, top, width):
    slide.shapes.add_picture(ICON_PATH, Inches(left), Inches(top), Inches(width))

# ============================================================
# SLIDE 1: Title / Primary Logo (Stacked)
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide1, ABYSS)

# Section label
add_text(slide1, 0, 0.4, 16, 0.5, "STRATA FORGE — BRAND IDENTITY", 11, CONCRETE, font_name='Calibri')

# Icon centered
add_icon(slide1, 6.0, 1.2, 4.0)

# STRATA
add_text(slide1, 2, 5.5, 12, 1.0, "STRATA", 72, INCANDESCENT, bold=True, spacing=600)

# FORGE with dashes
add_text(slide1, 2, 6.5, 12, 0.7, "—  FORGE  —", 36, ORANGE, spacing=800)

# Footer
add_text(slide1, 0, 8.3, 16, 0.4, "2026  ·  GAME STUDIO  ·  BRAND IDENTITY SYSTEM", 9, CONCRETE, font_name='Calibri')

# ============================================================
# SLIDE 2: Horizontal Layout
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide2, ABYSS)

add_text(slide2, 0, 0.4, 16, 0.5, "HORIZONTAL LAYOUT", 11, CONCRETE, font_name='Calibri')

# Icon left
add_icon(slide2, 3.5, 2.8, 2.5)

# STRATA right of icon
add_text(slide2, 6.5, 3.0, 7, 1.0, "STRATA", 60, INCANDESCENT, bold=True, alignment=PP_ALIGN.LEFT, spacing=500)

# FORGE
add_text(slide2, 6.5, 4.1, 7, 0.6, "—  FORGE  —", 28, ORANGE, alignment=PP_ALIGN.LEFT, spacing=700)

# ============================================================
# SLIDE 3: Background Variants (4 quadrants)
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide3, ABYSS)

add_text(slide3, 0, 0.4, 16, 0.5, "BACKGROUND VARIANTS", 11, CONCRETE, font_name='Calibri')

# 4 cards
cards = [
    {"label": "DARK", "bg": ABYSS, "border": TEAL, "text_color": INCANDESCENT, "forge_color": ORANGE},
    {"label": "LIGHT", "bg": LIGHT_BG, "border": RGBColor(0xDD, 0xDD, 0xDD), "text_color": RGBColor(0x1A, 0x1A, 0x1A), "forge_color": ORANGE_DEEP},
    {"label": "TEAL", "bg": TEAL_BG, "border": TEAL, "text_color": INCANDESCENT, "forge_color": ORANGE},
    {"label": "WARM", "bg": WARM_BG, "border": RGBColor(0x3A, 0x20, 0x15), "text_color": INCANDESCENT, "forge_color": ORANGE},
]

for i, card in enumerate(cards):
    col = i % 4
    x = 1.0 + col * 3.6
    y = 1.5

    # Background rectangle
    shape = slide3.shapes.add_shape(
        1,  # rectangle
        Inches(x), Inches(y), Inches(3.2), Inches(6.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = card["bg"]
    shape.line.color.rgb = card["border"]
    shape.line.width = Pt(1)

    # Icon
    add_icon(slide3, x + 0.85, y + 0.8, 1.5)

    # STRATA
    add_text(slide3, x, y + 3.2, 3.2, 0.6, "STRATA", 22, card["text_color"], bold=True, spacing=400)

    # FORGE
    add_text(slide3, x, y + 3.8, 3.2, 0.5, "FORGE", 13, card["forge_color"], spacing=600)

    # Label
    add_text(slide3, x, y + 5.6, 3.2, 0.4, card["label"], 9, CONCRETE, font_name='Calibri')

# ============================================================
# SLIDE 4: Icon Scaling
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide4, ABYSS)

add_text(slide4, 0, 0.4, 16, 0.5, "ICON SCALING", 11, CONCRETE, font_name='Calibri')

sizes = [
    {"label": "128px", "icon_w": 3.0, "x": 1.5},
    {"label": "64px", "icon_w": 1.5, "x": 5.5},
    {"label": "32px", "icon_w": 0.75, "x": 8.5},
    {"label": "16px", "icon_w": 0.4, "x": 11.0},
]

for s in sizes:
    # Dark frame
    frame_w = max(s["icon_w"] + 0.6, 1.2)
    frame_h = frame_w
    fx = s["x"]
    fy = 3.0

    shape = slide4.shapes.add_shape(1, Inches(fx), Inches(fy), Inches(frame_w), Inches(frame_h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ABYSS
    shape.line.color.rgb = TEAL
    shape.line.width = Pt(0.5)

    # Icon centered in frame
    ix = fx + (frame_w - s["icon_w"]) / 2
    iy = fy + (frame_h - s["icon_w"]) / 2
    add_icon(slide4, ix, iy, s["icon_w"])

    # Label
    add_text(slide4, fx, fy + frame_h + 0.2, frame_w, 0.3, s["label"], 11, CONCRETE, font_name='Calibri')

# ============================================================
# SLIDE 5: Brand Palette
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide5, ABYSS)

add_text(slide5, 0, 0.4, 16, 0.5, "BRAND PALETTE", 11, CONCRETE, font_name='Calibri')

palette = [
    {"name": "Teal", "hex": "#2A6B6B", "color": TEAL},
    {"name": "Forge Orange", "hex": "#E87830", "color": ORANGE},
    {"name": "Incandescent", "hex": "#FFF5E0", "color": WARM_WHITE},
    {"name": "Warm White", "hex": "#E8E0D0", "color": INCANDESCENT},
    {"name": "Abyss Black", "hex": "#0A0A12", "color": ABYSS},
]

for i, swatch in enumerate(palette):
    x = 2.0 + i * 2.5
    y = 2.5

    # Circle
    shape = slide5.shapes.add_shape(
        9,  # oval
        Inches(x), Inches(y), Inches(1.5), Inches(1.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = swatch["color"]
    if swatch["name"] == "Abyss Black":
        shape.line.color.rgb = TEAL
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()

    # Name
    add_text(slide5, x - 0.25, y + 1.8, 2.0, 0.4, swatch["name"], 12, INCANDESCENT)

    # Hex
    add_text(slide5, x - 0.25, y + 2.2, 2.0, 0.3, swatch["hex"], 11, CONCRETE, font_name='Calibri')

# ============================================================
# SLIDE 6: Usage Guidelines
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide6, ABYSS)

add_text(slide6, 0, 0.4, 16, 0.5, "USAGE GUIDELINES", 11, CONCRETE, font_name='Calibri')

guidelines = [
    {
        "title": "TYPOGRAPHY",
        "body": "STRATA: Cinzel Bold, spacing 0.28em\nFORGE: Cinzel Regular, spacing 0.55em\nColor: #E8E0D0 / #E87830"
    },
    {
        "title": "MINIMUM SIZE",
        "body": "Icon only: 16px minimum\nIcon + wordmark: 120px width min\nClear space: 1x icon width on all sides"
    },
    {
        "title": "DON'TS",
        "body": "No rotation or skewing\nNo color substitution\nNo drop shadow on light backgrounds\nNever on busy/patterned backgrounds"
    },
]

for i, g in enumerate(guidelines):
    x = 1.5 + i * 4.6
    y = 2.0

    # Card bg
    shape = slide6.shapes.add_shape(1, Inches(x), Inches(y), Inches(4.0), Inches(5.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0F, 0x10, 0x18)
    shape.line.color.rgb = RGBColor(0x1A, 0x2A, 0x2A)
    shape.line.width = Pt(0.5)

    # Title
    add_text(slide6, x + 0.3, y + 0.4, 3.4, 0.5, g["title"], 14, TEAL_LIGHT, bold=True, spacing=300)

    # Body
    txBox = slide6.shapes.add_textbox(Inches(x + 0.3), Inches(y + 1.2), Inches(3.4), Inches(3.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for line in g["body"].split("\n"):
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = CONCRETE
        p.font.name = "Calibri"
        p.space_after = Pt(6)

# Save
prs.save(OUTPUT_PATH)
print(f"PPTX saved: {OUTPUT_PATH}")
